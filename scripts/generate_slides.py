"""Generate the 1-hour PowerPoint presentation using CIMA template assets."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Paths ────────────────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parent.parent
BASI = ROOT / "POWERPOINT_CIMA" / "CIMA POWER POINT" / "BASI"
BG_WHITE = str(BASI / "SFONDO SLIDE 16-9-04.jpg")          # content
BG_BLUE_LOGO = str(BASI / "SFONDO SLIDE 16-9-03.jpg")      # section title
BG_BLUE_TITLE = str(BASI / "SFONDO SLIDE 16-9_Tavola disegno 1-01.jpg")  # title
BG_BLUE_CLOSE = str(BASI / "SFONDO SLIDE 16-9_Tavola disegno 1-02.jpg")  # closing
BG_BLUE_PLAIN = str(BASI / "SFONDO SLIDE 16-9 blu_Tavola disegno 1-05.jpg")  # plain blue
OUT = ROOT / "slides" / "DL_Best_Practices.pptx"

# ── Dimensions (13.33 × 7.5 in) ─────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Colors ───────────────────────────────────────────────────────────────
ORANGE = RGBColor(0xE8, 0x77, 0x00)
DARK_BLUE = RGBColor(0x00, 0x40, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x60, 0x60, 0x60)

# ── Helpers ──────────────────────────────────────────────────────────────

def _add_bg(slide, img_path: str):
    """Add a full-slide background image."""
    slide.shapes.add_picture(img_path, 0, 0, SLIDE_W, SLIDE_H)


def _add_textbox(slide, left, top, width, height, text, *,
                 font_name="Montserrat", font_size=Pt(20),
                 bold=False, color=DARK_BLUE, alignment=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = alignment
    except Exception:
        pass
    # Set vertical anchor
    txBox.text_frame._txBody.bodyPr.set("anchor", {
        MSO_ANCHOR.TOP: "t",
        MSO_ANCHOR.MIDDLE: "ctr",
        MSO_ANCHOR.BOTTOM: "b",
    }.get(anchor, "t"))
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def _add_bullet_slide(slide, title: str, bullets: list[str], *,
                      sub_bullets: dict[int, list[str]] | None = None):
    """Content slide with title + bullet list on white background."""
    _add_bg(slide, BG_WHITE)
    # Title
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 title, font_name="Montserrat SemiBold", font_size=Pt(32),
                 bold=True, color=DARK_BLUE)
    # Bullets
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = True

    if sub_bullets is None:
        sub_bullets = {}

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.name = "Montserrat"
        p.font.size = Pt(22)
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(8)
        p.level = 0
        from pptx.oxml.ns import qn
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(Inches(0.4)))
        pPr.set("indent", str(Inches(-0.3)))

        # Sub-bullets
        if i in sub_bullets:
            for sb in sub_bullets[i]:
                sp = tf.add_paragraph()
                sp.text = sb
                sp.font.name = "Montserrat"
                sp.font.size = Pt(18)
                sp.font.color.rgb = LIGHT_GRAY
                sp.space_after = Pt(4)
                sp.level = 1
                spPr = sp._p.get_or_add_pPr()
                spPr.set("marL", str(Inches(0.9)))
                spPr.set("indent", str(Inches(-0.25)))


def _section_slide(prs, title: str):
    """Dark blue section divider slide."""
    slide = prs.slides.add_slide(blank)
    _add_bg(slide, BG_BLUE_PLAIN)
    _add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(2.5),
                 title, font_name="Montserrat SemiBold", font_size=Pt(44),
                 bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    return slide


# ── Build presentation ───────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank = prs.slide_layouts[6]  # Blank layout

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
_add_bg(s, BG_BLUE_TITLE)
_add_textbox(s, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
             "Deep Learning", font_name="Montserrat ExtraBold",
             font_size=Pt(52), bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
_add_textbox(s, Inches(1.5), Inches(2.9), Inches(10), Inches(1.0),
             "Best Practices & Applications", font_name="Montserrat Light",
             font_size=Pt(36), color=ORANGE, alignment=PP_ALIGN.CENTER)
_add_textbox(s, Inches(1.5), Inches(4.3), Inches(10), Inches(0.6),
             "Module 2  |  CIMA Foundation PhD Course",
             font_name="Montserrat", font_size=Pt(20), color=WHITE,
             alignment=PP_ALIGN.CENTER)
_add_textbox(s, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
             "Luca Monaco  —  27 May 2026",
             font_name="Montserrat Light", font_size=Pt(18), color=WHITE,
             alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
_add_bullet_slide(s, "Agenda", [
    "Introduction to Deep Learning  (~10 min)",
    "The Modern ML Stack  (~10 min)",
    "Hands-on 1: Precipitation Post-Processing  (~15 min)",
    "Hands-on 2: Urban Thermal Comfort  (~15 min)",
    "Wrap-up & Q&A  (~10 min)",
])

# ═══════════════════════════════════════════════════════════════════════════
# SECTION 1 — Introduction to Deep Learning
# ═══════════════════════════════════════════════════════════════════════════
_section_slide(prs, "1. Introduction to\nDeep Learning")

# Slide: From Linear Models to Neural Networks
s = prs.slides.add_slide(blank)
_add_bullet_slide(s, "From Linear Models to Neural Networks", [
    "Linear regression: y = Wx + b",
    "Add non-linearity: activation functions (ReLU, sigmoid, tanh)",
    "Stack layers to learn hierarchical representations",
    "Universal approximation theorem: deep nets can model any function",
    'When to use DL vs. classical ML',
], sub_bullets={
    4: [
        "Large datasets with complex patterns \u2192 DL",
        "Small/tabular data, interpretability needed \u2192 classical ML",
    ],
})

# Slide: Key Building Blocks
s = prs.slides.add_slide(blank)
_add_bullet_slide(s, "Key Building Blocks", [
    "Neurons: weighted sum + activation",
    "Layers: Dense (fully-connected), Convolutional, Recurrent",
    "Loss functions: MSE, Cross-Entropy, Custom losses",
    "Optimizers: SGD, Adam, learning rate schedulers",
    "Regularization: Dropout, Batch Normalization, Weight Decay",
])

# Slide: Backpropagation & Gradient Descent
s = prs.slides.add_slide(blank)
_add_bullet_slide(s, "Backpropagation & Gradient Descent", [
    "Forward pass: compute predictions",
    "Loss: measure how wrong the predictions are",
    "Backward pass: chain rule propagates gradients layer by layer",
    "Parameter update: \u03B8 \u2190 \u03B8 \u2212 \u03B7 \u2207L(\u03B8)",
    "Mini-batch SGD: balance between speed and gradient noise",
    "Key challenges: vanishing/exploding gradients, local minima",
])

# Slide: Architecture Zoo
s = prs.slides.add_slide(blank)
_add_bullet_slide(s, "Architecture Overview", [
    "CNNs \u2014 Convolutional Neural Networks",
    "RNNs / LSTMs \u2014 Recurrent Neural Networks",
    "Transformers \u2014 Attention Is All You Need",
], sub_bullets={
    0: [
        "Spatial feature extraction via convolution & pooling",
        "Ideal for gridded geospatial data (precipitation fields, satellite imagery)",
    ],
    1: [
        "Sequential data with memory mechanisms",
        "Time series forecasting, language modeling",
    ],
    2: [
        "Self-attention: parallel processing of sequences",
        "Dominates NLP; increasingly used in weather/climate (ClimaX, Pangu-Weather)",
    ],
})

# ═══════════════════════════════════════════════════════════════════════════
# SECTION 2 — The Modern ML Stack
# ═══════════════════════════════════════════════════════════════════════════
_section_slide(prs, "2. The Modern ML Stack")

# Slide: Zarr + xarray
s = prs.slides.add_slide(blank)
_add_bullet_slide(s, "Zarr + xarray", [
    "Zarr: chunked, compressed, cloud-native N-dimensional arrays",
    "Why Zarr over NetCDF for ML?",
    "xarray integration: labeled dimensions, lazy loading",
    "ERA5 on Google Cloud: WeatherBench 2, ARCO-ERA5",
], sub_bullets={
    1: [
        "Parallel reads from cloud storage (GCS, S3)",
        "Efficient slicing \u2014 read only what you need",
        "No need to download entire files",
    ],
})

# Slide: PyTorch Lightning
s = prs.slides.add_slide(blank)
_add_bullet_slide(s, "PyTorch Lightning", [
    "Eliminates boilerplate: no manual .to(device), .zero_grad(), etc.",
    "LightningModule: define model, loss, optimizer in one class",
    "LightningDataModule: encapsulate data loading & splits",
    "Built-in callbacks: early stopping, checkpointing, logging",
    "Scales seamlessly: single GPU \u2192 multi-GPU \u2192 TPU",
])

# Slide: Hydra + MLFlow + uv
s = prs.slides.add_slide(blank)
_add_bullet_slide(s, "Hydra, MLFlow & uv", [
    "Hydra \u2014 Configuration Management",
    "MLFlow \u2014 Experiment Tracking",
    "uv \u2014 Modern Python Tooling",
], sub_bullets={
    0: [
        "YAML-based config composition & overrides",
        "Multirun sweeps for hyperparameter search",
    ],
    1: [
        "Log metrics, parameters, and artifacts",
        "Compare runs, reproduce results",
    ],
    2: [
        "Fast replacement for pip + virtualenv + pip-tools",
        "uv init / uv add / uv sync / uv run workflow",
    ],
})

# Slide: How It All Fits Together
s = prs.slides.add_slide(blank)
_add_bullet_slide(s, "How It All Fits Together", [
    "Data layer: Zarr store \u2192 xarray \u2192 PyTorch Dataset",
    "Config: Hydra YAML defines model, data, and training params",
    "Training: PyTorch Lightning Trainer handles the loop",
    "Tracking: MLFlow logs every experiment automatically",
    "Reproducibility: uv lock file + Hydra config = fully reproducible",
])

# ═══════════════════════════════════════════════════════════════════════════
# SECTION 3 — Hands-on 1
# ═══════════════════════════════════════════════════════════════════════════
_section_slide(prs, "3. Hands-on 1\nPrecipitation Post-Processing")

s = prs.slides.add_slide(blank)
_add_bullet_slide(s, "Hands-on 1: Precipitation Post-Processing", [
    "Goal: correct systematic biases in NWP precipitation forecasts",
    "Data: ERA5 reanalysis total precipitation (Zarr on Google Cloud)",
    "Model: simple CNN in PyTorch Lightning",
    "Workflow:",
], sub_bullets={
    3: [
        "1. Load data with xarray + Zarr",
        "2. Build CNN (conv layers \u2192 output field)",
        "3. Configure hyperparameters with Hydra (OmegaConf inline)",
        "4. Train & track with MLFlow",
        "5. Evaluate: bias maps, RMSE, visual comparison",
    ],
})

s = prs.slides.add_slide(blank)
_add_bullet_slide(s, "Notebook 1: Key Takeaways", [
    "Cloud-native data access removes the download bottleneck",
    "Lightning handles the training loop \u2014 focus on the science",
    "Hydra configs make experiments reproducible & sweepable",
    "Even simple CNNs can significantly reduce NWP bias",
    "Open the notebook \u2192 Google Colab link in the README",
])

# ═══════════════════════════════════════════════════════════════════════════
# SECTION 4 — Hands-on 2
# ═══════════════════════════════════════════════════════════════════════════
_section_slide(prs, "4. Hands-on 2\nUrban Thermal Comfort")

s = prs.slides.add_slide(blank)
_add_bullet_slide(s, "Hands-on 2: Urban Thermal Comfort", [
    "Goal: predict UTCI (thermal comfort index) from met + urban variables",
    "Data: ERA5-Land meteorological variables + land cover features (Zarr)",
    "Model: MLP (Multi-Layer Perceptron) in PyTorch Lightning",
    "Workflow:",
], sub_bullets={
    3: [
        "1. Load multi-variable urban climate dataset from Zarr",
        "2. Feature engineering & dataset preparation",
        "3. Build MLP with configurable hidden layers",
        "4. Train with Hydra config + MLFlow tracking",
        "5. Evaluate regression performance (RMSE, scatter plots)",
    ],
})

s = prs.slides.add_slide(blank)
_add_bullet_slide(s, "Notebook 2: Key Takeaways", [
    "Tabular geoscience data works well with simple MLPs",
    "Feature engineering matters \u2014 domain knowledge helps DL too",
    "Same stack (Lightning + Hydra + MLFlow) for different architectures",
    "Date-based train/val/test split avoids temporal data leakage",
    "Open the notebook \u2192 Google Colab link in the README",
])

# ═══════════════════════════════════════════════════════════════════════════
# SECTION 5 — Wrap-up
# ═══════════════════════════════════════════════════════════════════════════
_section_slide(prs, "5. Wrap-up & Q&A")

s = prs.slides.add_slide(blank)
_add_bullet_slide(s, "Key Takeaways", [
    "Deep Learning is a powerful tool \u2014 but not always the right one",
    "Modern stack (Zarr + Lightning + Hydra + MLFlow) makes DL accessible",
    "Cloud-native data formats unlock scalable, reproducible research",
    "Start simple, iterate fast, track everything",
])

s = prs.slides.add_slide(blank)
_add_bullet_slide(s, "Resources & Next Steps", [
    "PyTorch Lightning Docs \u2014 lightning.ai/docs/pytorch/stable/",
    "Zarr Docs \u2014 zarr.readthedocs.io",
    "Hydra Docs \u2014 hydra.cc/docs/intro/",
    "MLFlow Docs \u2014 mlflow.org/docs/latest/",
    "WeatherBench 2 \u2014 ERA5 Zarr on Google Cloud",
    "ARCO-ERA5 \u2014 Full ERA5 Zarr on Google Cloud",
    "Course repo \u2014 github.com/meteocima/Monaco_DLCourse",
])

# ═══════════════════════════════════════════════════════════════════════════
# CLOSING SLIDE
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
_add_bg(s, BG_BLUE_CLOSE)
_add_textbox(s, Inches(2), Inches(1.2), Inches(9), Inches(0.8),
             "Thank you!", font_name="Montserrat ExtraBold",
             font_size=Pt(44), bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
_add_textbox(s, Inches(2), Inches(2.0), Inches(9), Inches(0.6),
             "Questions?", font_name="Montserrat Light",
             font_size=Pt(28), color=ORANGE, alignment=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"Saved to {OUT}")
print(f"Total slides: {len(prs.slides)}")
